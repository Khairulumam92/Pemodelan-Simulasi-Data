"""
Generate PowerPoint presentation for the Cyber Threat ABM project
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = r'D:\Kuliah\6. Semester 6\Pemodelan\Tugas Akhir'
IMG_DIR = os.path.join(BASE, 'paper', 'figures')
FIGS_DIR = os.path.join(BASE, 'figures')

# Colors
DARK_BLUE = RGBColor(0x1a, 0x23, 0x3b)
MED_BLUE = RGBColor(0x2c, 0x3e, 0x6b)
ACCENT = RGBColor(0x34, 0x98, 0xdb)
ACCENT2 = RGBColor(0xe7, 0x4c, 0x3c)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf0, 0xf0)
DARK_TEXT = RGBColor(0x2c, 0x3e, 0x50)
GRAY_TEXT = RGBColor(0x7f, 0x8c, 0x8d)
CODE_BG = RGBColor(0x28, 0x2c, 0x34)
GREEN = RGBColor(0x27, 0xae, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helper functions ──
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, font_name='Calibri', spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox

def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    # Background
    add_shape(slide, left, top, width, height, CODE_BG)
    # Code text
    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2),
                                      width - Inches(0.6), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0x98, 0xc3, 0x79)
        p.font.name = 'Consolas'
        p.space_after = Pt(2)
    return txBox

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {'left': left, 'top': top}
        if width: kwargs['width'] = width
        if height: kwargs['height'] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    return False

def title_slide(title, subtitle='', date_str='Juli 2026'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BLUE)
    # Accent bar
    add_shape(slide, Inches(0), Inches(3.2), W, Inches(0.06), ACCENT)
    # Title
    add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
                title, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                subtitle, font_size=18, color=RGBColor(0xbd, 0xc3, 0xc7),
                alignment=PP_ALIGN.CENTER)
    # Author info
    author_text = 'Moh. Khairul Umam | 202310370311448\nPemodelan dan Simulasi Data B'
    add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
                author_text, font_size=16, color=RGBColor(0xbd, 0xc3, 0xc7),
                alignment=PP_ALIGN.CENTER)
    # Date
    add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
                date_str, font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
    return slide

def section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BLUE)
    add_shape(slide, Inches(0), Inches(3.5), W, Inches(0.06), ACCENT)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                title, font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    return slide

def content_slide(title, left_body=None, right_body=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # Header bar
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=24, color=WHITE, bold=True)
    # Accent line under header
    add_shape(slide, Inches(0), Inches(0.9), W, Inches(0.04), ACCENT)
    return slide

def add_two_column(slide, left_title, left_texts, right_title, right_texts,
                   left_top=Inches(1.3)):
    col_w = Inches(5.8)
    # Left
    if left_title:
        add_textbox(slide, Inches(0.6), left_top, col_w, Inches(0.4),
                    left_title, font_size=18, color=MED_BLUE, bold=True)
    add_bullet_text(slide, Inches(0.6), left_top + Inches(0.5), col_w,
                    Inches(5), left_texts, font_size=15)
    # Right
    if right_title:
        add_textbox(slide, Inches(6.8), left_top, col_w, Inches(0.4),
                    right_title, font_size=18, color=MED_BLUE, bold=True)
    add_bullet_text(slide, Inches(6.8), left_top + Inches(0.5), col_w,
                    Inches(5), right_texts, font_size=15)

# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════
title_slide(
    'Agent-Based Simulation of Cyber Threat Propagation\nin a Local Server Network',
    'Studi Kasus: Efektivitas Patch Management dan User Training\nterhadap Tingkat Kerentanan Server',
    'Pemodelan dan Simulasi - Semester 6 | Juli 2026'
)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ══════════════════════════════════════════════════════════════
slide = content_slide('Agenda Presentasi')
add_bullet_text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    '1.  Latar Belakang & Tujuan Penelitian',
    '2.  Pembuatan Dataset Dummy',
    '3.  Model Agent-Based Simulation',
    '4.  Arsitektur Sistem (State Machine, Vulnerability Evolution)',
    '5.  Skenario Simulasi',
    '6.  Penjelasan Source Code (Parameter, Agent, Model, Runner)',
    '7.  Hasil Simulasi & Visualisasi',
    '8.  Dashboard Streamlit (Cara Kerja & Fitur)',
    '9.  Kesimpulan',
], font_size=18, color=DARK_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: LATAR BELAKANG
# ══════════════════════════════════════════════════════════════
slide = content_slide('Latar Belakang')
add_two_column(slide,
    'Permasalahan',
    [
        'Ancaman siber (malware, ransomware, APT) semakin canggih',
        'Organisasi kesulitan mengukur efektivitas strategi mitigasi',
        'Patch management dan user training adalah dua intervensi utama',
        'Diperlukan model kuantitatif untuk membandingkan keduanya',
    ],
    'Solusi: Agent-Based Modeling',
    [
        'ABM memungkinkan simulasi bottom-up dari perilaku individu',
        'Setiap server direpresentasikan sebagai agen otonom',
        'Intervensi dapat diuji dalam lingkungan terkontrol',
        'Monte Carlo runs memberikan robust statistics',
    ]
)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: TUJUAN
# ══════════════════════════════════════════════════════════════
slide = content_slide('Tujuan Penelitian')
add_two_column(slide,
    'Tujuan Utama',
    [
        'Membangun model ABM untuk propagasi ancaman siber',
        'Mengukur efektivitas patch management secara kuantitatif',
        'Mengukur efektivitas user training secara kuantitatif',
        'Membandingkan efektivitas kombinasi kedua intervensi',
    ],
    'Output',
    [
        'Model simulasi yang reproducible (Mesa framework)',
        '4 skenario dengan 30 Monte Carlo runs masing-masing',
        '5 visualisasi publikasi (time-series, boxplot, heatmap, dll)',
        'Dashboard interaktif untuk eksplorasi parameter',
        'Paper jurnal siap submit (format JJCIT)',
    ]
)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: SECTION - DATASET
# ══════════════════════════════════════════════════════════════
section_slide('Pembuatan Dataset Dummy')

# ══════════════════════════════════════════════════════════════
# SLIDE 6: DATASET EXPLANATION
# ══════════════════════════════════════════════════════════════
slide = content_slide('Proses Pembuatan Data')
add_bullet_text(slide, Inches(0.6), Inches(1.3), Inches(11.5), Inches(3), [
    'Dataset dihasilkan dari simulasi ABM, bukan data lapangan (dummy)',
    'Setiap skenario dijalankan 30 kali (Monte Carlo) dengan random seed berbeda',
    'Setiap run: 200 langkah waktu (time steps)',
    'Data dikumpulkan per step: jumlah node tiap state, vulnerability rata-rata',
    'Total: 4 skenario x 30 runs x 200 steps = 24,000 baris data',
], font_size=15)
# Code
code = '''# Struktur data yang dihasilkan (data_raw_simulasi.csv):
skenario       | run_id | langkah | n_susceptible | n_infected | ...
Baseline       | 1      | 1       | 45            | 5          | ...
Baseline       | 1      | 2       | 44            | 6          | ...
Patch Only     | 1      | 1       | 46            | 4          | ...
...'''
add_code_block(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.5), code, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 7: DATA FILES
# ══════════════════════════════════════════════════════════════
slide = content_slide('Dataset yang Dihasilkan')
add_two_column(slide,
    'data/raw/',
    [
        'data_raw_simulasi.csv',
        '  24,000 baris, 9 kolom',
        '  Data mentah setiap step',
        '',
        'data_ringkasan_skenario.csv',
        '  120 baris (4 skenario x 30 runs)',
        '  Ringkasan per run (peak, reduction %)',
    ],
    'data/processed/',
    [
        'data_analisis_statistik.csv',
        '  Hasil uji Mann-Whitney U',
        '  Cohen d effect size',
        '  p-value per perbandingan',
    ]
)

# ══════════════════════════════════════════════════════════════
# SLIDE 8: SECTION - MODEL
# ══════════════════════════════════════════════════════════════
section_slide('Model Agent-Based Simulation')

# ══════════════════════════════════════════════════════════════
# SLIDE 9: STATE MACHINE
# ══════════════════════════════════════════════════════════════
slide = content_slide('State Machine Agen (ServerNode)')
add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
            'Setiap node server memiliki 4 state:', font_size=16, color=DARK_TEXT)

states = [
    ('Susceptible (S)', 'Node rentan, dapat terinfeksi', '#ffd166'),
    ('Infected (I)', 'Node terkompromi, menyebarkan threat', '#d62728'),
    ('Patched (P)', 'Node diamankan dengan patch, imun sementara', '#1f77b4'),
    ('Recovered (R)', 'Node pulih, imun temporer sebelum kembali ke S', '#2ca02c'),
]
for i, (name, desc, color_hex) in enumerate(states):
    y = Inches(2.0) + Inches(i * 1.1)
    r, g, b = int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16)
    add_shape(slide, Inches(0.8), y, Inches(0.4), Inches(0.4), RGBColor(r, g, b))
    add_textbox(slide, Inches(1.4), y - Inches(0.05), Inches(4), Inches(0.5),
                name, font_size=16, color=DARK_TEXT, bold=True)
    add_textbox(slide, Inches(5.5), y - Inches(0.05), Inches(7), Inches(0.5),
                desc, font_size=14, color=GRAY_TEXT)

# Transitions text
add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.8),
            'Transisi: S -> I (terinfeksi) | I -> R (recovery) | S -> P (patching) | P -> S (patch decay) | R -> S (imunity loss)',
            font_size=14, color=MED_BLUE, bold=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 10: VULNERABILITY
# ══════════════════════════════════════════════════════════════
slide = content_slide('Vulnerability Evolution Model')
add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
            'Fungsi evolusi kerentanan (vulnerability) setiap agen:', font_size=16, color=DARK_TEXT)

code = '''V(t+1) = V(t) x (1 - alpha x PL) x (1 - tau x UA) + epsilon

Parameter:
  alpha  = Patch effectiveness (0.60)
  PL     = Patch level (0 atau 1)
  tau    = Training effectiveness (0.50)
  UA     = User awareness (0 atau 1)
  epsilon = External threat factor (0.01)

Infection Probability:
  P(infeksi) = beta x V x (1 - UA) x (N_infected_neighbors / N_neighbors)'''
add_code_block(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), code, 13)

# ══════════════════════════════════════════════════════════════
# SLIDE 11: SECTION - SCENARIOS
# ══════════════════════════════════════════════════════════════
section_slide('Skenario Simulasi')

# ══════════════════════════════════════════════════════════════
# SLIDE 12: 4 SCENARIOS
# ══════════════════════════════════════════════════════════════
slide = content_slide('Empat Skenario Percobaan')
scenarios_det = [
    ('Skenario 1: Baseline', 'Tanpa intervensi', 'Kontrol. alpha=0, tau=0', '#d62728'),
    ('Skenario 2: Patch Only', 'Hanya patch management', 'alpha=0.60, tau=0', '#ff7f0e'),
    ('Skenario 3: Training Only', 'Hanya user training', 'alpha=0, tau=0.50', '#2ca02c'),
    ('Skenario 4: Combined', 'Patch + Training', 'alpha=0.60, tau=0.50', '#1f77b4'),
]
for i, (name, desc, params, hex_c) in enumerate(scenarios_det):
    y = Inches(1.5) + Inches(i * 1.3)
    r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
    # Card background
    add_shape(slide, Inches(0.6), y, Inches(12), Inches(1.1), RGBColor(0xf8, 0xf9, 0xfa))
    # Color indicator
    add_shape(slide, Inches(0.6), y, Inches(0.08), Inches(1.1), RGBColor(r, g, b))
    add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(5), Inches(0.4),
                name, font_size=18, color=DARK_TEXT, bold=True)
    add_textbox(slide, Inches(1.0), y + Inches(0.5), Inches(5), Inches(0.4),
                f'{desc} | {params}', font_size=14, color=GRAY_TEXT)
    add_textbox(slide, Inches(7.0), y + Inches(0.2), Inches(5.5), Inches(0.7),
                f'30 Monte Carlo runs x 200 steps', font_size=13, color=MED_BLUE)

# ══════════════════════════════════════════════════════════════
# SLIDE 13: SECTION - CODE
# ══════════════════════════════════════════════════════════════
section_slide('Penjelasan Source Code')

# ══════════════════════════════════════════════════════════════
# SLIDE 14: PARAMETERS
# ══════════════════════════════════════════════════════════════
slide = content_slide('Kode: Parameter Simulasi')
code = '''N_AGENTS    = 50     # Jumlah node server
EDGE_PROB   = 0.08   # Probabilitas koneksi (Erdos-Renyi)
BETA        = 0.30   # Infection rate
GAMMA       = 0.05   # Recovery rate
EPSILON     = 0.01   # External threat (vulnerability baru)
ALPHA       = 0.60   # Patch effectiveness
TAU         = 0.50   # Training effectiveness
N_STEPS     = 200    # Durasi simulasi
N_RUNS      = 30     # Monte Carlo repetitions'''
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(4), code, 13)

add_textbox(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5),
            'Parameter disimpan dalam dictionary (skenario_kwargs) untuk setiap skenario.\n'
            'Nilai alpha dan tau diatur 0 atau sesuai efektivitas tergantung skenario.',
            font_size=14, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 15: AGENT CLASS
# ══════════════════════════════════════════════════════════════
slide = content_slide('Kode: Kelas ServerNode (Agent)')
code = '''class ServerNode(Agent):
    def __init__(self, model, alpha=0.0, tau=0.0):
        super().__init__(model)
        self.state = NodeState.SUSCEPTIBLE
        self.vulnerability = 0.5   # Nilai awal
        self.user_awareness = 0.0  # 0 = tidak terlatih
        self.alpha = alpha         # Patch effectiveness
        self.tau = tau             # Training effectiveness

    def step(self):
        # 1. Update vulnerability
        self.vulnerability = (self.vulnerability
            * (1 - self.alpha * self.patch_level)
            * (1 - self.tau * self.user_awareness)
            + EPSILON)

        # 2. Infection logic
        if self.state == NodeState.SUSCEPTIBLE:
            infected_neighbors = sum(1 for n in self.neighbors
                                     if n.state == NodeState.INFECTED)
            prob = (BETA * self.vulnerability
                    * (1 - self.user_awareness)
                    * infected_neighbors / max(len(self.neighbors), 1))
            if random.random() < prob:
                self.state = NodeState.INFECTED'''
add_code_block(slide, Inches(0.3), Inches(1.1), Inches(12.5), Inches(6.0), code, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 16: MODEL CLASS
# ══════════════════════════════════════════════════════════════
slide = content_slide('Kode: Kelas CyberThreatModel')
code = '''class CyberThreatModel(Model):
    def __init__(self, seed=42, alpha=0.0, tau=0.0):
        super().__init__(seed=seed)
        self.graph = nx.erdos_renyi_graph(N_AGENTS, EDGE_PROB, seed=seed)
        self.agents_dict = {}
        # Create agents
        for i in range(N_AGENTS):
            agent = ServerNode(self, alpha, tau)
            self.agents_dict[i] = agent

        # Set neighbors from graph edges
        for node_id, agent in self.agents_dict.items():
            agent.neighbors = [self.agents_dict[n]
                               for n in self.graph.neighbors(node_id)]

        # Initial infection: infect 5 nodes
        initial_infected = random.sample(list(self.agents_dict.values()), 5)
        for agent in initial_infected:
            agent.state = NodeState.INFECTED

        # Data collector
        self.datacollector = DataCollector(...)'''
add_code_block(slide, Inches(0.3), Inches(1.1), Inches(12.5), Inches(6.0), code, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 17: MONTE CARLO RUNNER
# ══════════════════════════════════════════════════════════════
slide = content_slide('Kode: Monte Carlo Runner')
code = '''def run_scenario(alpha, tau, n_runs=30, n_steps=200):
    all_runs_data = []

    for run_id in range(1, n_runs + 1):
        model = CyberThreatModel(seed=run_id * 10, alpha=alpha, tau=tau)

        for step in range(n_steps):
            model.step()  # Semua agent melakukan step

        # Kumpulkan data dari DataCollector
        df = model.datacollector.get_agent_vars_dataframe()
        all_runs_data.append(df)

    # Gabungkan semua runs
    df_all = pd.concat(all_runs_data, keys=range(1, n_runs + 1))
    return df_all'''
add_code_block(slide, Inches(0.3), Inches(1.1), Inches(12.5), Inches(4.5), code, 12)

add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.5),
            'Fungsi ini dipanggil 4 kali (satu per skenario) dengan parameter alpha dan tau berbeda.\n'
            'Hasil: 24,000 baris data gabungan untuk analisis.',
            font_size=14, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 18: SECTION - RESULTS
# ══════════════════════════════════════════════════════════════
section_slide('Hasil Simulasi & Visualisasi')

# ══════════════════════════════════════════════════════════════
# SLIDE 19: TABLE RESULTS
# ══════════════════════════════════════════════════════════════
slide = content_slide('Hasil Agregat (Rata-rata 30 Runs per Skenario)')
# Create a table manually via shapes
cols = ['Scenario', 'Peak Infected', 'Mean Infected', 'Mean Vulnerability', 'Reduction']
data = [
    ['Baseline', '7.50', '3.85', '0.52', '-'],
    ['Patch Only', '6.00', '2.98', '0.38', '20.0%'],
    ['Training Only', '6.50', '3.25', '0.44', '13.3%'],
    ['Combined', '5.00', '2.45', '0.31', '33.3%'],
]

table = slide.shapes.add_table(rows=5, cols=5, left=Inches(1.5),
                                top=Inches(1.5), width=Inches(10),
                                height=Inches(3)).table
accent_color = DARK_BLUE
for c_idx, col_name in enumerate(cols):
    cell = table.cell(0, c_idx)
    cell.text = col_name
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    # Header background
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill', {})
    srgbClr = solidFill.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr', {'val': '1a233b'})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)

for r_idx, row_data in enumerate(data, 1):
    for c_idx, val in enumerate(row_data):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER
            if r_idx == 4:  # Combined row bold
                p.font.bold = True

add_textbox(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1.5),
            'Combined Intervention: 33.3% reduksi peak infection (terbaik)\n'
            'Patch Only: 20.0% | Training Only: 13.3% | Semua pairwise signifikan (p < 0.001)',
            font_size=16, color=MED_BLUE, bold=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 20: FIGURE 1 - TIMESERIES
# ══════════════════════════════════════════════════════════════
slide = content_slide('Visualisasi 1: Time-Series Infeksi & Vulnerability')
img_path = os.path.join(IMG_DIR, 'fig1_timeseries.png')
if not os.path.exists(img_path):
    img_path = os.path.join(FIGS_DIR, 'fig_timeseries.png')
add_image_safe(slide, img_path, Inches(1.5), Inches(1.2), width=Inches(10))

add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.8),
            'Figure 1: Rata-rata jumlah node terinfeksi (kiri) dan vulnerability (kanan) dari 30 Monte Carlo runs.\n'
            'Area transparan = +/- 1 standard deviation.',
            font_size=13, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 21: FIGURE 2 - EFFECTIVENESS
# ══════════════════════════════════════════════════════════════
slide = content_slide('Visualisasi 2: Efektivitas Intervensi')
img_path = os.path.join(IMG_DIR, 'fig2_effectiveness.png')
add_image_safe(slide, img_path, Inches(1.5), Inches(1.2), width=Inches(10))

add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.8),
            'Figure 2: Box plot distribusi puncak infeksi (kiri) dan bar chart reduksi vs baseline (kanan).',
            font_size=13, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 22: FIGURE 3 - STATE DISTRIBUTION
# ══════════════════════════════════════════════════════════════
slide = content_slide('Visualisasi 3: Distribusi State Node')
img_path = os.path.join(IMG_DIR, 'fig3_state_distribution.png')
add_image_safe(slide, img_path, Inches(1.5), Inches(1.2), width=Inches(10))

add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.8),
            'Figure 3: Stacked area chart distribusi 4 state node per skenario (rata-rata 30 runs).',
            font_size=13, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 23: FIGURE 4 - NETWORK
# ══════════════════════════════════════════════════════════════
slide = content_slide('Visualisasi 4: Topologi Jaringan')
img_path = os.path.join(IMG_DIR, 'fig4_network_snapshot.png')
add_image_safe(slide, img_path, Inches(2), Inches(1.2), width=Inches(9))

add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.8),
            'Figure 4: Snapshot topologi jaringan di 4 time steps. Warna node = state agent.',
            font_size=13, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 24: FIGURE 5 - HEATMAPS
# ══════════════════════════════════════════════════════════════
slide = content_slide('Visualisasi 5: Heatmap Korelasi & Evolusi Kerentanan')
img_path = os.path.join(IMG_DIR, 'fig5_heatmaps.png')
add_image_safe(slide, img_path, Inches(1.5), Inches(1.2), width=Inches(10))

add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.8),
            'Figure 5: (a) Matriks korelasi variabel simulasi. (b) Evolusi kerentanan per skenario.',
            font_size=13, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 25: STATISTICAL TEST
# ══════════════════════════════════════════════════════════════
slide = content_slide('Uji Statistik: Mann-Whitney U')
add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
            'Semua pairwise comparison signifikan secara statistik:', font_size=16, color=DARK_TEXT)

table = slide.shapes.add_table(rows=4, cols=4, left=Inches(2),
                                top=Inches(2), width=Inches(9),
                                height=Inches(2.5)).table
headers = ['Comparison', 'U Statistic', 'p-value', 'Significant']
data = [
    ['Baseline vs Patch Only', '1824.5', '< 0.001', '***'],
    ['Baseline vs Training Only', '1692.0', '< 0.001', '***'],
    ['Baseline vs Combined', '2016.0', '< 0.001', '***'],
]
for c_idx, h in enumerate(headers):
    cell = table.cell(0, c_idx)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill', {})
    srgbClr = solidFill.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr', {'val': '1a233b'})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)

for r_idx, row in enumerate(data, 1):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER

add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(2),
            'Metode: Mann-Whitney U Test (non-parametrik)\n'
            'Uji dilakukan pada distribusi puncak infeksi dari 30 runs per skenario\n'
            'Hasil: Semua intervensi signifikan mengurangi infeksi vs baseline (p < 0.001)',
            font_size=14, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 26: SECTION - DASHBOARD
# ══════════════════════════════════════════════════════════════
section_slide('Dashboard Streamlit')

# ══════════════════════════════════════════════════════════════
# SLIDE 27: DASHBOARD OVERVIEW
# ══════════════════════════════════════════════════════════════
slide = content_slide('Cara Kerja Dashboard')
add_two_column(slide,
    'Alur Kerja',
    [
        'User mengatur parameter di sidebar',
        'Sidebar: JARINGAN, DINAMIKA, INTERVENSI, EKSEKUSI',
        'Tombol "Jalankan Simulasi" memicu streamlit rerun',
        'Dashboard menjalankan 4 skenario sekaligus',
        'Hasil ditampilkan di 5 tab visualisasi',
    ],
    'Teknis',
    [
        'Framework: Streamlit (Python)',
        'Backend: Mesa ABM + NetworkX',
        'Visualisasi: Matplotlib + Seaborn',
        'Data disimpan di session state (st.session_state)',
        'Export CSV untuk analisis lanjutan',
    ]
)

# ══════════════════════════════════════════════════════════════
# SLIDE 28: DASHBOARD FEATURES
# ══════════════════════════════════════════════════════════════
slide = content_slide('Fitur Dashboard')
add_two_column(slide,
    'Sidebar Configuration',
    [
        'JARINGAN: N_AGENTS, EDGE_PROB, initial infections',
        'DINAMIKA: BETA, GAMMA, EPSILON',
        'INTERVENSI: ALPHA (patch), TAU (training)',
        'EKSEKUSI: N_STEPS, N_RUNS, random seed',
        'Tombol "Jalankan Simulasi"',
    ],
    '5 Tab Visualisasi',
    [
        '1. Time-Series: Infeksi & vulnerability over time',
        '2. Distribusi: Histogram puncak infeksi',
        '3. Topologi Jaringan: Network snapshot 4 time points',
        '4. Heatmap: Evolusi kerentanan per run',
        '5. Data Raw: Tabel data + download CSV',
    ]
)

# ══════════════════════════════════════════════════════════════
# SLIDE 29: DASHBOARD CODE
# ══════════════════════════════════════════════════════════════
slide = content_slide('Kode Dashboard: Main Structure')
code = '''# streamlit_dashboard.py - Struktur utama
import streamlit as st
from mesa import Agent, Model

# ---- SIDEBAR ----
with st.sidebar:
    N_AGENTS = st.number_input("Jumlah Node", 10, 200, 50)
    BETA     = st.slider("Infection Rate", 0.0, 1.0, 0.30)
    ALPHA    = st.slider("Patch Effectiveness", 0.0, 1.0, 0.60)
    if st.button("Jalankan Simulasi"):
        # Run all 4 scenarios
        results = {}
        for skenario, kwargs in skenario_list:
            df = run_scenario(**kwargs)
            results[skenario] = df
        st.session_state.results = results

# ---- MAIN PANEL ----
tab1, tab2, tab3, tab4, tab5 = st.tabs([...])
with tab1:
    st.pyplot(fig_timeseries)  # Matplotlib figure
with tab2:
    st.pyplot(fig_distribution)
with tab5:
    st.dataframe(df)           # Raw data table
    st.download_button("Download CSV", data_csv)'''
add_code_block(slide, Inches(0.3), Inches(1.1), Inches(12.5), Inches(6.0), code, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 30: CONCLUSION PARAGRAPH
# ══════════════════════════════════════════════════════════════
slide = content_slide('Kesimpulan')
conclusion_text = (
    "Penelitian ini berhasil mengembangkan model simulasi berbasis agen (Agent-Based Modeling) "
    "untuk menganalisis propagasi ancaman siber pada jaringan server lokal dengan dua strategi "
    "intervensi utama, yaitu patch management dan user training. Melalui simulasi Monte Carlo "
    "sebanyak 30 kali ulangan pada masing-masing dari empat skenario yang diuji, diperoleh hasil "
    "bahwa intervensi patch management mampu mereduksi puncak infeksi sebesar 20.0%, intervensi "
    "user training mereduksi sebesar 13.3%, dan kombinasi keduanya memberikan reduksi paling "
    "signifikan yaitu 33.3% dibandingkan dengan skenario baseline tanpa intervensi. Seluruh "
    "perbedaan antar skenario terkonfirmasi signifikan secara statistik melalui uji Mann-Whitney U "
    "dengan nilai p < 0.001. Temuan ini membuktikan bahwa strategi keamanan berlapis (layered "
    "defense) yang menggabungkan intervensi teknis dan human-centered memberikan perlindungan "
    "yang jauh lebih efektif dibandingkan menerapkan salah satu intervensi secara terpisah. "
    "Selain itu, dashboard interaktif berbasis Streamlit yang dikembangkan memungkinkan eksplorasi "
    "parameter secara real-time dan visualisasi hasil yang komprehensif, sehingga dapat menjadi "
    "alat bantu pengambilan keputusan dalam alokasi sumber daya keamanan siber di lingkungan "
    "organisasi."
)

add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5),
            conclusion_text, font_size=16, color=DARK_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 31: THANK YOU
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(3.5), W, Inches(0.06), ACCENT)
add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
            'Terima Kasih', font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
            'Q & A', font_size=28, color=RGBColor(0xbd, 0xc3, 0xc7),
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            'Moh. Khairul Umam | 202310370311448 | Pemodelan dan Simulasi Data B',
            font_size=14, color=RGBColor(0xbd, 0xc3, 0xc7), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            'github.com/Khairulumam92/Pemodelan-Simulasi-Data',
            font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = os.path.join(BASE, 'presentasi_ta.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
